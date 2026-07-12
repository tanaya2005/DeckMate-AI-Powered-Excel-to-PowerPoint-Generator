"""
ppt_builder.py — JSON slide plan + real DataFrames → .pptx (Day 3)

Takes the structured slide plan from llm_planner.py and the REAL dataframes
from data_parser.py, then builds a native, editable PowerPoint deck using
python-pptx.

ALL chart numbers come from pandas — the LLM only decides structure and
writes narrative text.

Pipeline position:
    Excel → [data_parser] → summary → [llm_planner] → slide plan →
    [ppt_builder] → downloadable .pptx
"""

import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Brand / Theme constants
# ---------------------------------------------------------------------------

# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colour palette
COLOR_PRIMARY    = RGBColor(0x1B, 0x2A, 0x4A)  # Dark navy
COLOR_SECONDARY  = RGBColor(0x2E, 0x86, 0xAB)  # Teal
COLOR_ACCENT     = RGBColor(0xE8, 0x6F, 0x51)  # Coral
COLOR_BG_LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)  # Off-white
COLOR_TEXT_DARK   = RGBColor(0x33, 0x33, 0x33)  # Near-black
COLOR_TEXT_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)  # White
COLOR_SUBTITLE    = RGBColor(0x66, 0x66, 0x66)  # Grey

# Branded 6-Color Palette
COLOR_TEAL   = RGBColor(0x2E, 0x86, 0xAB)  # Teal
COLOR_CORAL  = RGBColor(0xE8, 0x6F, 0x51)  # Coral
COLOR_AMBER  = RGBColor(0xF2, 0xB7, 0x05)  # Amber
COLOR_SAGE   = RGBColor(0x76, 0x9A, 0x6B)  # Sage Green
COLOR_SLATE  = RGBColor(0x6C, 0x7A, 0x89)  # Slate
COLOR_LAVENDER = RGBColor(0x9B, 0x59, 0xB6) # Purple/Lavender

CHART_COLORS = [
    COLOR_TEAL,
    COLOR_CORAL,
    COLOR_AMBER,
    COLOR_SAGE,
    COLOR_SLATE,
    COLOR_LAVENDER
]

# Chart type mapping
CHART_TYPE_MAP = {
    "bar":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "pie":  XL_CHART_TYPE.PIE,
    "line": XL_CHART_TYPE.LINE_MARKERS,
}

# Layout positions (Wider elements to fill up whitespace)
TITLE_LEFT   = Inches(0.6)
TITLE_TOP    = Inches(0.4)
TITLE_WIDTH  = Inches(12.133)
TITLE_HEIGHT = Inches(0.8)

# Lowered chart/insight box to accommodate KPI callouts on top
CHART_LEFT   = Inches(0.6)
CHART_TOP    = Inches(2.5)       # Shifted down from 1.4 to 2.5
CHART_WIDTH  = Inches(7.8)
CHART_HEIGHT = Inches(4.1)       # Shrunk height from 5.0 to 4.1

INSIGHT_LEFT   = Inches(8.7)
INSIGHT_TOP    = Inches(2.5)       # Shifted down from 1.4 to 2.5
INSIGHT_WIDTH  = Inches(4.0)
INSIGHT_HEIGHT = Inches(4.1)       # Shrunk height from 5.0 to 4.1

TABLE_LEFT   = Inches(0.6)
TABLE_TOP    = Inches(2.5)
TABLE_WIDTH  = Inches(12.133)

FOOTER_LEFT   = Inches(0.6)
FOOTER_TOP    = Inches(6.9)
FOOTER_WIDTH  = Inches(12.133)
FOOTER_HEIGHT = Inches(0.4)


# ---------------------------------------------------------------------------
# Column name resolution (trust boundary)
# ---------------------------------------------------------------------------

def _normalize(name: str) -> str:
    """Normalize a column name for comparison: lowercase, strip whitespace."""
    return "".join(name.lower().split())


def _resolve_columns(
    requested_cols: List[str],
    df: pd.DataFrame,
) -> Tuple[List[str], List[str]]:
    """
    Resolve LLM-provided column names against the actual DataFrame columns.

    Uses case-insensitive, whitespace-tolerant matching.

    Returns
    -------
    (matched_cols, unmatched_cols)
        matched_cols  — list of ACTUAL DataFrame column names (in the order requested)
        unmatched_cols — list of LLM column names that couldn't be matched
    """
    actual_cols = list(df.columns)
    # Build a lookup: normalized_name -> actual_name
    lookup = {}
    for col in actual_cols:
        norm = _normalize(col)
        if norm not in lookup:  # first match wins (handles duplicates)
            lookup[norm] = col

    matched = []
    unmatched = []

    for req in requested_cols:
        norm_req = _normalize(req)
        if norm_req in lookup:
            matched.append(lookup[norm_req])
        else:
            unmatched.append(req)

    return matched, unmatched


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def compute_kpi(df, kpi_spec):
    """
    kpi_spec examples:
    {"type": "extreme_category", "group_by": "Category", "value_col": "Current Stock", "agg": "sum", "which": "min"}
    {"type": "avg_per_group", "group_by": "Supplier", "value_col": "Current Stock", "agg": "sum"}
    {"type": "filtered_rate", "filter_col": "Reorder Status", "filter_contains": "Reorder"}
    {"type": "filtered_sum", "filter_col": "Reorder Status", "filter_contains": "Reorder", "value_col": "Current Stock"}
    {"type": "peak_point", "x_col": "Order Date", "value_col": "Total Quantity Ordered"}
    """
    if kpi_spec["type"] == "extreme_category":
        grouped = df.groupby(kpi_spec["group_by"])[kpi_spec["value_col"]].agg(kpi_spec["agg"])
        idx = grouped.idxmin() if kpi_spec["which"] == "min" else grouped.idxmax()
        return idx, grouped[idx]

    if kpi_spec["type"] == "avg_per_group":
        grouped = df.groupby(kpi_spec["group_by"])[kpi_spec["value_col"]].agg(kpi_spec["agg"])
        return grouped.mean()  # mean of the GROUPED sums, never df[value_col].mean() directly

    if kpi_spec["type"] == "filtered_rate":
        mask = df[kpi_spec["filter_col"]].str.contains(kpi_spec["filter_contains"], case=False, na=False)
        return (mask.sum() / len(df)) * 100

    if kpi_spec["type"] == "filtered_sum":
        mask = df[kpi_spec["filter_col"]].str.contains(kpi_spec["filter_contains"], case=False, na=False)
        return df.loc[mask, kpi_spec["value_col"]].sum()

    if kpi_spec["type"] == "peak_point":
        idx = df[kpi_spec["value_col"]].idxmax()
        return df.loc[idx, kpi_spec["x_col"]], df.loc[idx, kpi_spec["value_col"]]


def _resolve_kpi_spec_cols(df: pd.DataFrame, spec: dict) -> dict:
    resolved_spec = spec.copy()
    for col_key in ["group_by", "value_col", "filter_col", "x_col"]:
        if col_key in spec and spec[col_key]:
            resolved, _ = _resolve_columns([spec[col_key]], df)
            if resolved:
                resolved_spec[col_key] = resolved[0]
    return resolved_spec


def _format_single_value(val) -> str:
    if isinstance(val, (float, np.floating)):
        if float(val).is_integer():
            return f"{int(val):,}"
        else:
            return f"{val:,.1f}"
    elif isinstance(val, (int, np.integer)):
        return f"{int(val):,}"
    elif hasattr(val, "strftime"):
        return val.strftime("%Y-%m-%d")
    else:
        return str(val)


def _format_kpi_value(val, spec_type=None) -> str:
    if isinstance(val, tuple):
        key_part, val_part = val
        if hasattr(key_part, "strftime"):
            key_part = key_part.strftime("%Y-%m-%d")
        else:
            key_part = str(key_part)
        return f"{key_part} ({_format_single_value(val_part)})"
    else:
        formatted = _format_single_value(val)
        if spec_type == "filtered_rate":
            formatted += "%"
        return formatted


def _compute_kpis(
    kpi_defs: list,
    df: pd.DataFrame,
    matched_cols: List[str],
) -> Dict[str, str]:
    """
    Compute real KPI values from the DataFrame using the LLM's spec definitions.
    """
    if not kpi_defs or df is None or df.empty:
        return {}

    result: Dict[str, str] = {}
    for spec in kpi_defs[:3]:  # max 3 KPIs
        label = spec.get("label", "")
        if not label:
            continue
        try:
            resolved_spec = _resolve_kpi_spec_cols(df, spec)
            val = compute_kpi(df, resolved_spec)
            formatted = _format_kpi_value(val, resolved_spec.get("type"))
            
            print(f"[KPI Computation] Label: '{label}' | Spec: {spec} -> Value: '{formatted}'")
            result[label] = formatted
        except Exception as e:
            print(f"[KPI Computation] ERROR computing KPI '{label}' with spec {spec}: {e}")
            continue

    return result


def _process_insight_text(
    insight_text: str,
    df: pd.DataFrame,
    matched_cols: List[str],
    slide_title: str = "",
) -> str:
    """
    Calculate top-line numeric facts from the real dataframe and inject them
    into the insight text by replacing placeholders.
    """
    if not insight_text:
        return ""
    if df is None or df.empty or not matched_cols:
        return insight_text

    try:
        # Determine category vs value columns
        category_col = matched_cols[0]
        value_cols = matched_cols[1:] if len(matched_cols) > 1 else []

        val_col = None
        if value_cols:
            # Try to find the first numeric value column
            for col in value_cols:
                if pd.api.types.is_numeric_dtype(df[col]):
                    val_col = col
                    break
            # Fallback to the first value column if none are numeric
            if val_col is None:
                val_col = value_cols[0]

        if not val_col:
            # Count occurrences of category_col
            agg_df = df[category_col].value_counts().reset_index()
            agg_df.columns = [category_col, "Count"]
            val_col = "Count"
        else:
            # Aggregate: group by category, sum numeric value column
            if pd.api.types.is_numeric_dtype(df[val_col]):
                agg_df = df.groupby(category_col)[val_col].sum().reset_index()
            else:
                agg_df = df.groupby(category_col)[val_col].count().reset_index()

        agg_df = agg_df.sort_values(by=val_col, ascending=False).reset_index(drop=True)

        total_val = agg_df[val_col].sum()
        mean_val = agg_df[val_col].mean()
        count_val = len(agg_df)

        top_cat, top_val, top_pct = "N/A", "0", "0.0%"
        sec_cat, sec_val, sec_pct = "N/A", "0", "0.0%"
        bot_cat, bot_val = "N/A", "0"

        if count_val > 0:
            row = agg_df.iloc[0]
            top_cat = str(row[category_col])
            t_val = row[val_col]
            top_val = f"{t_val:,.2f}" if isinstance(t_val, float) else f"{t_val:,}"
            pct = (t_val / total_val * 100) if total_val else 0.0
            top_pct = f"{pct:.1f}%"

            bot_row = agg_df.iloc[-1]
            bot_cat = str(bot_row[category_col])
            b_val = bot_row[val_col]
            bottom_val = f"{b_val:,.2f}" if isinstance(b_val, float) else f"{b_val:,}"

        if count_val > 1:
            row = agg_df.iloc[1]
            sec_cat = str(row[category_col])
            s_val = row[val_col]
            sec_val = f"{s_val:,.2f}" if isinstance(s_val, float) else f"{s_val:,}"
            pct = (s_val / total_val * 100) if total_val else 0.0
            sec_pct = f"{pct:.1f}%"

        total_val_str = f"{total_val:,.2f}" if isinstance(total_val, float) else f"{total_val:,}"
        mean_val_str = f"{mean_val:,.2f}"
        count_str = f"{count_val:,}"

        replacements = {
            "{{TOP_CATEGORY}}": top_cat,
            "{{TOP_VALUE}}": top_val,
            "{{TOP_PCT}}": top_pct,
            "{{SECOND_CATEGORY}}": sec_cat,
            "{{SECOND_VALUE}}": sec_val,
            "{{SECOND_PCT}}": sec_pct,
            "{{BOTTOM_CATEGORY}}": bot_cat,
            "{{BOTTOM_VALUE}}": bottom_val,
            "{{TOTAL_VALUE}}": total_val_str,
            "{{MEAN_VALUE}}": mean_val_str,
            "{{COUNT}}": count_str,
        }

        processed = insight_text
        for placeholder, replacement in replacements.items():
            processed = processed.replace(placeholder, replacement)

        print(f"[Insight Replacements] Slide: '{slide_title}' -> Replacements: {replacements}")
        return processed

    except Exception as e:
        print(f"[Insight Replacements] ERROR processing slide '{slide_title}': {e}")
        return insight_text


def _add_kpi_callouts(slide, kpis: dict):
    """Draw a row of styled KPI callout boxes at the top of the slide."""
    if not kpis:
        return

    # Filter out empty entries
    valid_kpis = {k: v for k, v in kpis.items() if k and v}
    if not valid_kpis:
        return

    # Limit to max 3 KPIs to fit the width
    items = list(valid_kpis.items())[:3]
    count = len(items)

    # Box coordinates
    kpi_y = Inches(1.3)
    kpi_h = Inches(0.9)
    total_w = Inches(12.133)
    
    # Calculate box width and spacing dynamically
    box_w = Inches(3.6)
    spacing = Inches(0.4)
    
    # Left margin offset to center the KPIs container
    container_w = (box_w * count) + (spacing * (count - 1))
    start_x = Inches(0.6) + (total_w - container_w) / 2

    from pptx.enum.shapes import MSO_SHAPE
    
    for i, (label, val) in enumerate(items):
        box_x = start_x + i * (box_w + spacing)
        
        # Add a rounded rectangle background card for each KPI
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, kpi_y, box_w, kpi_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_SECONDARY
        card.line.width = Pt(1.5)
        
        # Add textbox inside the card for label and value
        txBox = slide.shapes.add_textbox(box_x, kpi_y + Inches(0.05), box_w, kpi_h - Inches(0.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # KPI Value (large, bold)
        p_val = tf.paragraphs[0]
        p_val.text = str(val)
        p_val.font.size = Pt(20)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_PRIMARY
        p_val.alignment = PP_ALIGN.CENTER
        
        # KPI Label (small, gray)
        p_lbl = tf.add_paragraph()
        p_lbl.text = str(label).upper()
        p_lbl.font.size = Pt(9)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_SECONDARY
        p_lbl.alignment = PP_ALIGN.CENTER

def _add_title_textbox(slide, title_text: str):
    """Add a styled title text box to the top of a slide, stripping markdown bold."""
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    parts = title_text.split("**")
    for idx, part in enumerate(parts):
        if not part:
            continue
        run = p.add_run()
        run.text = part
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY


def _add_insight_textbox(slide, insight_text: str, left=None, top=None,
                         width=None, height=None):
    """Add a styled insight text box, rendering subheadings and bullet points properly over a light card background."""
    x = left or INSIGHT_LEFT
    y = top or INSIGHT_TOP
    w = width or INSIGHT_WIDTH
    h = height or INSIGHT_HEIGHT

    # Add light background card shape (rounded rectangle or rectangle)
    from pptx.enum.shapes import MSO_SHAPE
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_BG_LIGHT
    card.line.fill.background() # No border
    
    # Add text box on top of the card (with slight margins)
    margin = Inches(0.25)
    txBox = slide.shapes.add_textbox(
        x + margin,
        y + margin,
        w - (margin * 2),
        h - (margin * 2),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Title label
    label_p = tf.paragraphs[0]
    label_p.text = "ANALYSIS & KEY INSIGHTS"
    label_p.font.size = Pt(11)
    label_p.font.bold = True
    label_p.font.color.rgb = COLOR_SECONDARY
    label_p.space_after = Pt(12)

    # Split the insight text into segments
    lines = insight_text.replace("\\n", "\n").split("\n")
    for line in lines:
        line_str = line.strip()
        if not line_str:
            continue

        p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = Pt(18)

        # Check if it's a bullet point
        if line_str.startswith("•") or line_str.startswith("-") or line_str.startswith("*"):
            # Strip the bullet char
            content = line_str.lstrip("•-* ").strip()
            # Style bullet paragraph
            p.level = 0
            p.space_before = Pt(4)
            # Add bullet symbol explicitly
            run_bullet = p.add_run()
            run_bullet.text = "•  "
            run_bullet.font.size = Pt(13)
            run_bullet.font.bold = True
            run_bullet.font.color.rgb = COLOR_SECONDARY

            # Parse bold text within the bullet
            _add_formatted_runs(p, content)
        else:
            # Normal line or heading
            _add_formatted_runs(p, line_str)


def _add_formatted_runs(paragraph, text: str):
    """Parse text and add bold/normal runs (looks for **text**)."""
    parts = text.split("**")
    for idx, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_TEXT_DARK

        # Odd indices are inside the ** ** block, so they should be bolded
        if idx % 2 == 1:
            run.font.bold = True
            run.font.color.rgb = COLOR_PRIMARY
        else:
            run.font.bold = False


def _add_footer(slide, slide_num: int, total_slides: int):
    """Add a subtle footer with slide number."""
    txBox = slide.shapes.add_textbox(FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH, FOOTER_HEIGHT)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"DeckMate  |  Slide {slide_num} of {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_SUBTITLE
    p.alignment = PP_ALIGN.RIGHT


def _style_chart(chart, chart_type: str):
    """Apply consistent styling to a chart."""
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)

    # Style the series/points colours
    has_points = (chart_type in ("pie", "bar"))
    
    for i, series in enumerate(chart.series):
        # Determine series level color
        color = CHART_COLORS[i % len(CHART_COLORS)]
        
        # If single series and is bar or pie, color individual points differently
        if len(chart.series) == 1 and has_points and hasattr(series, "points"):
            for p_idx in range(len(series.points)):
                pt = series.points[p_idx]
                pt_color = CHART_COLORS[p_idx % len(CHART_COLORS)]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = pt_color
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

        # Line charts: style the line and markers
        if chart_type == "line":
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2.5)
            if hasattr(series, 'marker'):
                series.marker.style = 8  # circle
                series.marker.size = 8

    # Value axis formatting (not applicable to pie)
    if chart_type != "pie":
        try:
            value_axis = chart.value_axis
            value_axis.has_title = False
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            value_axis.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            value_axis.tick_labels.font.size = Pt(9)
            value_axis.tick_labels.font.color.rgb = COLOR_SUBTITLE

            category_axis = chart.category_axis
            category_axis.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            category_axis.tick_labels.font.size = Pt(9)
            category_axis.tick_labels.font.color.rgb = COLOR_SUBTITLE
        except Exception:
            pass  # some chart types may not have these axes


def _prepare_chart_data(
    df: pd.DataFrame,
    category_col: str,
    value_cols: List[str],
    chart_type: str,
) -> CategoryChartData:
    """
    Prepare chart data from the REAL dataframe.

    For bar/pie: aggregate by category column, sum the value columns.
    For line: use category as-is (typically dates/quarters), values as series.
    """
    chart_data = CategoryChartData()

    if chart_type in ("bar", "pie"):
        # Aggregate: group by category, sum numeric value columns
        agg_dict = {}
        for vc in value_cols:
            if pd.api.types.is_numeric_dtype(df[vc]):
                agg_dict[vc] = "sum"
            else:
                agg_dict[vc] = "count"

        if agg_dict:
            grouped = df.groupby(category_col, sort=True).agg(agg_dict).reset_index()
        else:
            # All value columns are non-numeric — count occurrences
            grouped = df[category_col].value_counts().reset_index()
            grouped.columns = [category_col, "Count"]
            value_cols = ["Count"]

    elif chart_type == "line":
        # For line charts, sort by the category (usually date/time)
        grouped = df.copy()
        try:
            grouped = grouped.sort_values(category_col)
        except Exception:
            pass

        # If there are many rows, aggregate by category
        if len(grouped) > 20:
            agg_dict = {vc: "mean" for vc in value_cols if pd.api.types.is_numeric_dtype(df[vc])}
            if agg_dict:
                grouped = grouped.groupby(category_col).agg(agg_dict).reset_index()
    else:
        grouped = df.copy()

    # Set categories
    categories = [str(v) for v in grouped[category_col].tolist()]
    chart_data.categories = categories

    # Add series
    for vc in value_cols:
        if vc in grouped.columns:
            values = grouped[vc].tolist()
            # Ensure numeric
            values = [float(v) if pd.notna(v) else 0.0 for v in values]
            chart_data.add_series(vc, values)

    return chart_data


def _build_chart_slide(
    prs: Presentation,
    slide_info: dict,
    df: pd.DataFrame,
    matched_cols: List[str],
    slide_num: int,
    total_slides: int,
):
    """Build a slide with a native chart + insight text with dynamic height budgeting."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    chart_type = slide_info["chart_type"]
    title = slide_info.get("title", "")
    insight = slide_info.get("insight_text", "")

    # Add title
    _add_title_textbox(slide, title)

    # Add KPI Callouts (computed from REAL data, not LLM values)
    kpi_defs = slide_info.get("kpis", [])
    computed_kpis = _compute_kpis(kpi_defs, df, matched_cols)
    _add_kpi_callouts(slide, computed_kpis)

    # Determine category vs value columns
    category_col = matched_cols[0]
    value_cols = matched_cols[1:] if len(matched_cols) > 1 else []

    if not value_cols:
        if pd.api.types.is_numeric_dtype(df[category_col]):
            value_cols = [category_col]
            category_col = df.index.name or "Index"
            if category_col not in df.columns:
                df = df.reset_index()
        else:
            counts = df[category_col].value_counts().reset_index()
            counts.columns = [category_col, "Count"]
            df = counts
            value_cols = ["Count"]

    # Dynamic Height Budgeting:
    # Slide height = 7.5. Footer top = 6.9. Title & KPIs take up top 2.3.
    # Total remaining space between y = 2.4 and y = 6.8 is 4.4 inches.
    content_top = Inches(2.4)
    max_content_h = Inches(4.3)
    
    # We fit the chart into the max height
    chart_h = max_content_h

    # Prepare chart data from REAL dataframe
    chart_data = _prepare_chart_data(df, category_col, value_cols, chart_type)
    xl_chart_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Add chart
    chart_frame = slide.shapes.add_chart(
        xl_chart_type,
        CHART_LEFT, content_top, CHART_WIDTH, chart_h,
        chart_data,
    )

    # Style the chart
    _style_chart(chart_frame.chart, chart_type)

    # Add insight text box (right side)
    _add_insight_textbox(slide, insight, left=INSIGHT_LEFT, top=content_top, width=INSIGHT_WIDTH, height=chart_h)

    # Add footer
    _add_footer(slide, slide_num, total_slides)


def _build_table_slide(
    prs: Presentation,
    slide_info: dict,
    df: pd.DataFrame,
    matched_cols: List[str],
    slide_num: int,
    total_slides: int,
):
    """Build a slide with a data table + insight text with dynamic row shrinking and height budgeting."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    title = slide_info.get("title", "")
    insight = slide_info.get("insight_text", "")

    _add_title_textbox(slide, title)

    # Add KPI Callouts (computed from REAL data, not LLM values)
    kpi_defs = slide_info.get("kpis", [])
    computed_kpis = _compute_kpis(kpi_defs, df, matched_cols)
    _add_kpi_callouts(slide, computed_kpis)

    # Prepare table data (no hard truncate, show all rows requested by plan)
    table_df = df[matched_cols].copy()
    rows = len(table_df) + 1  # +1 for header
    cols = len(matched_cols)

    # Dynamic Height Budgeting:
    # Slide height = 7.5. Available vertical area for content: from y = 2.4 to y = 6.8 (4.4 inches).
    content_top = Inches(2.4)
    max_height_budget = Inches(4.3)
    
    # Reserve space for insights (minimum guaranteed height)
    min_insight_h = Inches(1.3)
    max_table_h = max_height_budget - min_insight_h - Inches(0.2) # leaves buffer

    # Calculate row height dynamically
    default_row_h = Inches(0.4)
    needed_height = default_row_h * rows
    
    font_size = 10
    row_h = default_row_h

    if needed_height > max_table_h:
        # Scale down row height to fit the budget
        row_h = max_table_h / rows
        # Sensible minimum row height (approx 0.2 inches)
        min_row_h = Inches(0.22)
        if row_h < min_row_h:
            print(f"[ppt_builder] WARNING: Table has too many rows ({rows}). Even at min row height, it may overflow.")
            row_h = min_row_h
        
        # Scale font down proportionally (min 8pt, default 10pt)
        ratio = float(row_h) / float(default_row_h)
        font_size = max(8, int(10 * ratio))

    table_height = row_h * rows
    
    # Scale column widths if there are many columns
    default_col_w = TABLE_WIDTH / cols
    col_w = default_col_w
    col_font_size = font_size
    if cols > 6:
        col_font_size = max(8, font_size - 1)

    table_shape = slide.shapes.add_table(
        rows, cols,
        TABLE_LEFT, content_top, TABLE_WIDTH, table_height,
    )
    table = table_shape.table

    # Set row heights explicitly
    for r_idx in range(rows):
        table.rows[r_idx].height = row_h

    # Style header row
    for j, col_name in enumerate(matched_cols):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.font.size = Pt(max(9, col_font_size + 1))
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    # Fill data rows
    for i in range(len(table_df)):
        for j, col_name in enumerate(matched_cols):
            cell = table.cell(i + 1, j)
            val = table_df.iloc[i][col_name]
            if pd.isna(val):
                cell.text = "—"
            elif isinstance(val, float):
                cell.text = f"{val:,.2f}"
            else:
                cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(col_font_size)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.alignment = PP_ALIGN.LEFT

            # Alternate row shading
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BG_LIGHT

    # Position the insight box dynamically BELOW the table
    insight_top = content_top + table_height + Inches(0.15)
    insight_h = max(min_insight_h, max_height_budget - table_height)
    
    _add_insight_textbox(
        slide, insight,
        left=TABLE_LEFT, top=insight_top,
        width=TABLE_WIDTH, height=insight_h,
    )

    _add_footer(slide, slide_num, total_slides)


def _build_text_slide(
    prs: Presentation,
    slide_info: dict,
    slide_num: int,
    total_slides: int,
    is_fallback: bool = False,
):
    """Build a text-only slide (title, takeaways, or fallback for unmatched columns)."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    title = slide_info.get("title", "")
    insight = slide_info.get("insight_text", "")

    # Centre the title for text slides
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.2),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    parts_title = title.split("**")
    for idx, part in enumerate(parts_title):
        if not part:
            continue
        run = p.add_run()
        run.text = part
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY

    # Body text
    if insight:
        body_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5), Inches(10.0), Inches(3.0),
        )
        bf = body_box.text_frame
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp.line_spacing = Pt(28)
        
        parts_body = insight.replace("\\n", "\n").split("**")
        for idx, part in enumerate(parts_body):
            if not part:
                continue
            run = bp.add_run()
            run.text = part
            run.font.size = Pt(18)
            run.font.color.rgb = COLOR_TEXT_DARK
            if idx % 2 == 1:
                run.font.bold = True
                run.font.color.rgb = COLOR_PRIMARY

    # Fallback notice
    if is_fallback:
        warn_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(6.0), Inches(10.0), Inches(0.5),
        )
        wf = warn_box.text_frame
        wp = wf.paragraphs[0]
        wp.text = "(Chart could not be rendered — column mismatch with data source)"
        wp.font.size = Pt(10)
        wp.font.color.rgb = COLOR_SUBTITLE
        wp.font.italic = True
        wp.alignment = PP_ALIGN.CENTER

    _add_footer(slide, slide_num, total_slides)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_presentation(
    slide_plan: dict,
    dataframes: dict,
    output_path: str,
) -> str:
    """
    Build a PowerPoint presentation from the slide plan and real DataFrames.
    Automatically prepends a Title Slide and appends a Closing Slide.

    Parameters
    ----------
    slide_plan : dict
        Output from llm_planner.generate_slide_plan(). Must have a "slides" key.
    dataframes : dict
        Dict of {sheet_name: pd.DataFrame} from data_parser.parse_excel().
    output_path : str
        Where to save the .pptx file.

    Returns
    -------
    str — the output path.
    """
    prs = Presentation()

    # Set 16:9 slide size
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = slide_plan.get("slides", [])
    
    # We will insert a Title Slide (Slide 1) and a Closing Slide (Last Slide)
    # Total count = len(slides) + 2
    total_slides = len(slides) + 2

    # 1. PREPEND TITLE SLIDE
    title_layout = prs.slide_layouts[6] # blank
    title_slide = prs.slides.add_slide(title_layout)
    
    # Draw dark navy accent panel at the bottom or left
    from pptx.enum.shapes import MSO_SHAPE
    accent_bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.8), Inches(13.333), Inches(0.4)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_SECONDARY
    accent_bar.line.fill.background()
    
    # Title Text Box
    t_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.5))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_plan.get("title_slide_title", "Data Analysis & Business Report")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Subtitle Text Box
    s_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.333), Inches(1.0))
    s_tf = s_box.text_frame
    s_tf.word_wrap = True
    sp = s_tf.paragraphs[0]
    from datetime import datetime
    current_date = datetime.now().strftime("%B %d, %Y")
    sp.text = f"Prepared automatically on {current_date} | Powered by DeckMate"
    sp.font.size = Pt(16)
    sp.font.color.rgb = COLOR_SUBTITLE
    
    _add_footer(title_slide, 1, total_slides)

    # 2. RENDER CONTENT SLIDES
    for idx, slide_info in enumerate(slides):
        slide_num = idx + 2 # +1 for index, +1 for Title Slide offset
        chart_type = slide_info.get("chart_type", "text").lower().strip()
        source_sheet = slide_info.get("source_sheet", "none")
        data_columns = slide_info.get("data_columns", [])

        # Find the source dataframe
        df = None
        if source_sheet.lower() != "none":
            df = dataframes.get(source_sheet)

        matched_cols = []
        if df is not None and not df.empty and data_columns:
            # ------- COLUMN VALIDATION (trust boundary) -------
            matched_cols, unmatched_cols = _resolve_columns(data_columns, df)

            if unmatched_cols:
                print(f"[ppt_builder] WARNING: Slide {slide_num} ('{slide_info.get('title', '')}') "
                      f"— LLM requested columns {unmatched_cols} which don't exist in "
                      f"sheet '{source_sheet}'. Available: {list(df.columns)}")

            if matched_cols:
                # Compute and process insight text placeholders
                raw_insight = slide_info.get("insight_text", "")
                slide_info["insight_text"] = _process_insight_text(
                    raw_insight, df, matched_cols, slide_info.get("title", "")
                )

        # ------- TEXT SLIDES (no data needed or requested as text slide) -------
        if chart_type == "text" or source_sheet.lower() == "none" or df is None or df.empty or len(matched_cols) < 1:
            is_fallback = (source_sheet.lower() != "none" and (df is None or df.empty or len(matched_cols) < 1))
            if is_fallback:
                print(f"[ppt_builder] WARNING: Falling back to text slide for slide {slide_num} due to missing data/columns.")
            _build_text_slide(prs, slide_info, slide_num, total_slides, is_fallback=is_fallback)
            continue

        # ------- BUILD THE APPROPRIATE SLIDE TYPE -------
        start_slide_count = len(prs.slides)
        try:
            df_copy = df.copy()
            if chart_type in ("bar", "pie", "line"):
                _build_chart_slide(
                    prs, slide_info, df_copy, matched_cols, slide_num, total_slides,
                )
            elif chart_type == "table":
                _build_table_slide(
                    prs, slide_info, df_copy, matched_cols, slide_num, total_slides,
                )
            else:
                _build_text_slide(prs, slide_info, slide_num, total_slides)

        except Exception as e:
            # If slide building failed, clean up any slide(s) added during this failed attempt
            while len(prs.slides) > start_slide_count:
                slide_id_list = prs.slides._sldIdLst
                del slide_id_list[-1]

            # If chart building fails for any reason, fallback to text
            print(f"[ppt_builder] ERROR building slide {slide_num}: {e}. "
                  f"Falling back to text slide.")
            _build_text_slide(prs, slide_info, slide_num, total_slides, is_fallback=True)

    # 3. APPEND CLOSING SLIDE (Thank You Slide)
    closing_layout = prs.slide_layouts[6]
    closing_slide = prs.slides.add_slide(closing_layout)
    
    # Draw navy backdrop block on the left
    left_block = closing_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5)
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = COLOR_PRIMARY
    left_block.line.fill.background()
    
    # Add Large Thank You Text
    close_box = closing_slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(7.0), Inches(2.5))
    close_tf = close_box.text_frame
    close_tf.word_wrap = True
    
    p_thank = close_tf.paragraphs[0]
    p_thank.text = "Thank You"
    p_thank.font.size = Pt(54)
    p_thank.font.bold = True
    p_thank.font.color.rgb = COLOR_PRIMARY
    
    p_sub = close_tf.add_paragraph()
    p_sub.text = "Questions & Discussion"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_SECONDARY
    p_sub.space_before = Pt(12)
    
    _add_footer(closing_slide, total_slides, total_slides)

    # Validate that total slide count matches expectations
    assert len(prs.slides) == total_slides, f"Mismatch: {len(prs.slides)} slides built but total_slides was {total_slides}"

    # Scan every text frame & table cell in the entire presentation for literal substring "{{"
    for slide_idx, sld in enumerate(prs.slides):
        for shape in sld.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if "{{" in txt:
                    print(f"[ValidationError] Unfilled placeholder found in Slide {slide_idx+1}: {txt}")
                    raise ValueError(f"Unfilled template placeholder found in Slide {slide_idx+1}")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if "{{" in cell.text:
                            print(f"[ValidationError] Unfilled placeholder in table cell in Slide {slide_idx+1}: {cell.text}")
                            raise ValueError(f"Unfilled template placeholder found in table in Slide {slide_idx+1}")

    # Save
    prs.save(output_path)
    print(f"[ppt_builder] Saved presentation to: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# CLI test block
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import json
    import numpy as np

    print("=" * 60)
    print("DAY 3 TEST: PPTX Builder")
    print("=" * 60)

    # Sanity checks using the real excels/inventory_management.xlsx
    from pathlib import Path
    excel_file = Path("excels/inventory_management.xlsx")
    if excel_file.exists():
        df_stock = pd.read_excel(excel_file, sheet_name="Stock Levels")
        df_pending = pd.read_excel(excel_file, sheet_name="Pending Orders")
        
        grouped_category_sum = df_stock.groupby("Category")["Current Stock"].agg("sum")
        supplier_sum = df_stock.groupby("Supplier")["Current Stock"].agg("sum")
        reorder_mask = df_stock["Reorder Status"].str.contains("Reorder", case=False, na=False)
        
        # Sanity checks
        assert grouped_category_sum.idxmin() == "Biography"  # lowest category
        assert grouped_category_sum.idxmax() == "Self-Help"   # highest category
        assert round(supplier_sum.mean(), 1) == 829.5          # avg stock per supplier
        assert round((reorder_mask.sum()/len(df_stock))*100, 1) == 12.0  # reorder rate
        
        # Convert date to datetime for timestamp assertion comparison
        df_pending_dt = df_pending.copy()
        df_pending_dt["Order Date"] = pd.to_datetime(df_pending_dt["Order Date"])
        assert df_pending_dt.loc[df_pending_dt["Total Quantity Ordered"].idxmax(), "Order Date"] == pd.Timestamp("2026-06-20")  # peak day
        
        print("All KPI sanity checks passed.")
    else:
        print("[WARNING] excels/inventory_management.xlsx not found, skipping sanity check assertions.")

    bookstore_file = Path("excels/bookstore_billing.xlsx")
    if bookstore_file.exists():
        df_books = pd.read_excel(bookstore_file, sheet_name="Billing Transactions")
        
        # 1. Top genre == "Sci-Fi" with count 24
        genre_counts = df_books.groupby("Genre")["Quantity"].count()
        assert genre_counts.idxmax() == "Sci-Fi"
        assert genre_counts.max() == 24
        
        # 2. Top payment method == "UPI" with count 27
        pay_counts = df_books.groupby("Payment Method")["Quantity"].count()
        assert pay_counts.idxmax() == "UPI"
        assert pay_counts.max() == 27
        
        # 3. Average quantity ≈ 2.67
        avg_qty = df_books["Quantity"].mean()
        assert abs(avg_qty - 2.67) < 0.01
        
        # 4. Average price per unit ≈ 13.36
        avg_price = df_books["Price Per Unit"].mean()
        assert abs(avg_price - 13.36) < 0.01
        
        print("All bookstore sanity checks passed.")
    else:
        print("[WARNING] excels/bookstore_billing.xlsx not found, skipping bookstore checks.")

    # Build a fake slide plan + dataframes to test each slide type
    test_df_sales = pd.DataFrame({
        "Region": ["North", "South", "East", "West"] * 3,
        "Product": ["Widget A", "Widget B", "Widget C",
                     "Widget A", "Widget B", "Widget C",
                     "Widget A", "Widget B", "Widget C",
                     "Widget A", "Widget B", "Widget C"],
        "Revenue": [120000, 95000, 145000, 88000,
                    110000, 72000, 130000, 98000,
                    155000, 91000, 105000, 85000],
        "Units Sold": [450, 380, 520, 310,
                       410, 290, 480, 370,
                       560, 340, 400, 320],
    })

    test_df_marketing = pd.DataFrame({
        "Quarter": ["Q1", "Q1", "Q2", "Q2", "Q3", "Q3", "Q4", "Q4"],
        "Channel": ["Online", "Retail", "Online", "Retail",
                     "Online", "Retail", "Online", "Retail"],
        "Spend": [50000, 35000, 62000, 41000, 71000, 45000, 80000, 52000],
        "ROI %": [12.5, 8.3, 15.1, 9.7, 18.2, 11.0, 20.5, 13.4],
    })

    dataframes = {
        "Regional Sales": test_df_sales,
        "Marketing Spend": test_df_marketing,
    }

    slide_plan = {
        "slides": [
            {
                "title": "Sales Performance Overview",
                "source_sheet": "none",
                "chart_type": "text",
                "data_columns": [],
                "insight_text": "This deck presents an overview of regional sales performance and marketing ROI for the fiscal year, prepared for senior leadership review.",
                "kpis": [],
            },
            {
                "title": "Revenue by Region",
                "source_sheet": "Regional Sales",
                "chart_type": "bar",
                "data_columns": ["Region", "Revenue"],
                "insight_text": "The {{TOP_CATEGORY}} region led revenue generation with {{TOP_VALUE}} total, representing {{TOP_PCT}} of the {{TOTAL_VALUE}} total revenue. The second region is {{SECOND_CATEGORY}} with {{SECOND_VALUE}} ({{SECOND_PCT}}).",
                "kpis": [
                    {"label": "Total Revenue", "type": "filtered_sum", "filter_col": "Region", "filter_contains": "", "value_col": "Revenue"},
                    {"label": "Top Region", "type": "extreme_category", "group_by": "Region", "value_col": "Revenue", "agg": "sum", "which": "max"},
                ],
            },
            {
                "title": "Product Revenue Share",
                "source_sheet": "Regional Sales",
                "chart_type": "pie",
                "data_columns": ["Product", "Revenue"],
                "insight_text": "{{TOP_CATEGORY}} commands the largest share of total revenue at {{TOP_PCT}} ({{TOP_VALUE}}), followed by {{SECOND_CATEGORY}} at {{SECOND_PCT}}.",
                "kpis": [
                    {"label": "Total Revenue", "type": "filtered_sum", "filter_col": "Product", "filter_contains": "", "value_col": "Revenue"},
                    {"label": "Top Product", "type": "extreme_category", "group_by": "Product", "value_col": "Revenue", "agg": "sum", "which": "max"},
                ],
            },
            {
                "title": "Marketing ROI Trend",
                "source_sheet": "Marketing Spend",
                "chart_type": "line",
                "data_columns": ["Quarter", "ROI %"],
                "insight_text": "Average ROI across {{COUNT}} periods is {{MEAN_VALUE}}%, with a maximum of {{TOP_VALUE}}% from the leading quarter. Lowest was {{BOTTOM_VALUE}}% from {{BOTTOM_CATEGORY}}.",
                "kpis": [
                    {"label": "Avg ROI", "type": "avg_per_group", "group_by": "Quarter", "value_col": "ROI %", "agg": "mean"},
                    {"label": "Max ROI", "type": "filtered_sum", "filter_col": "Quarter", "filter_contains": "", "value_col": "ROI %"}, # using filtered_sum as max placeholder for demo
                ],
            },
            {
                "title": "Marketing Spend Breakdown",
                "source_sheet": "Marketing Spend",
                "chart_type": "table",
                "data_columns": ["Quarter", "Channel", "Spend", "ROI %"],
                "insight_text": "A total of {{TOTAL_VALUE}} was spent across the channels, with {{TOP_CATEGORY}} having the highest spend at {{TOP_VALUE}} representing {{TOP_PCT}} of total spend.",
                "kpis": [
                    {"label": "Total Spend", "type": "filtered_sum", "filter_col": "Quarter", "filter_contains": "", "value_col": "Spend"},
                    {"label": "Top Spend Channel", "type": "extreme_category", "group_by": "Channel", "value_col": "Spend", "agg": "sum", "which": "max"},
                ],
            },
            {
                # This slide deliberately has a WRONG column name to test
                # the fuzzy-match fallback
                "title": "Hallucinated Column Test",
                "source_sheet": "Regional Sales",
                "chart_type": "bar",
                "data_columns": ["Region", "TotalRevenue"],  # "TotalRevenue" doesn't exist
                "insight_text": "This slide tests the fallback when the LLM hallucinates a column name that doesn't exist in the data.",
                "kpis": [
                    {"label": "Total Revenue", "type": "filtered_sum", "filter_col": "Region", "filter_contains": "", "value_col": "TotalRevenue"},
                ],
            },
            {
                # Test fuzzy match: "ROI%" (no space) vs actual "ROI %"
                "title": "Fuzzy Column Match Test",
                "source_sheet": "Marketing Spend",
                "chart_type": "bar",
                "data_columns": ["Channel", "ROI%"],  # "ROI%" should match "ROI %"
                "insight_text": "This tests that 'ROI%' fuzzy-matches to the actual column 'ROI %' despite missing whitespace.",
                "kpis": [
                    {"label": "Avg ROI", "type": "avg_per_group", "group_by": "Channel", "value_col": "ROI%", "agg": "mean"},
                ],
            },
            {
                "title": "Key Takeaways",
                "source_sheet": "none",
                "chart_type": "text",
                "data_columns": [],
                "insight_text": "East region is the clear revenue leader. Marketing ROI is trending upward. Online channels significantly outperform Retail on ROI. Recommend increased online investment.",
                "kpis": [],
            },
        ]
    }

    output_file = "test_output.pptx"

    print(f"\nBuilding {len(slide_plan['slides'])} slides...\n")
    build_presentation(slide_plan, dataframes, output_file)

    print(f"\nDone! Open '{output_file}' to verify:")
    print("  - Slide 1: Text (title/overview)")
    print("  - Slide 2: Bar chart (Revenue by Region)")
    print("  - Slide 3: Pie chart (Product Revenue Share)")
    print("  - Slide 4: Line chart (Marketing ROI Trend)")
    print("  - Slide 5: Table (Marketing Spend Breakdown)")
    print("  - Slide 6: Text FALLBACK (hallucinated column 'TotalRevenue')")
    print("  - Slide 7: Bar chart (fuzzy match 'ROI%' -> 'ROI %')")
    print("  - Slide 8: Text (Key Takeaways)")
